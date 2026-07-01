"""
Process knowledge base: extract text from PPTX files and scraped website data,
then save as structured JSON chunks ready for vector embedding.

Usage:
    python process_knowledge.py

Output:
    backend/data/chunks/knowledge_chunks.json
"""

import os
import json
from pptx import Presentation

# ── Paths ─────────────────────────────────────────────────────────────────────
_DATA_DIR = os.path.dirname(__file__)
KNOWLEDGE_BASE_DIR = os.path.join(_DATA_DIR, "..", "..", "knowledge-base")
SCRAPED_DIR = os.path.join(KNOWLEDGE_BASE_DIR, "scraped_data")
CHUNKS_DIR = os.path.join(_DATA_DIR, "chunks")
OUTPUT_FILE = os.path.join(CHUNKS_DIR, "knowledge_chunks.json")

PPTX_FILES = {
    "hi": "Fan Basics - Hindi.pptx",
    "ta": "Fan Basics - Tamil.pptx",
    "ml": "Fan Basics - Malayalam.pptx",
}

WEBSITE_CHUNK_WORDS = 250  # words per website chunk


# ── PPTX extraction ───────────────────────────────────────────────────────────

def extract_text_from_pptx(path: str, language: str) -> list[dict]:
    """Extract text from each slide of a PPTX file as individual chunks."""
    chunks = []
    if not os.path.exists(path):
        print(f"  ⚠️  Not found: {path}")
        return chunks

    prs = Presentation(path)
    filename = os.path.basename(path)

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        if texts:
            chunks.append({
                "text": " ".join(texts),
                "source": filename,
                "source_type": "pptx",
                "language": language,
                "slide": slide_num,
            })

    print(f"  ✅ {filename}: {len(chunks)} slides extracted")
    return chunks


# ── Website data extraction ───────────────────────────────────────────────────

def extract_text_from_scraped_website(path: str) -> list[dict]:
    """Read scraped website txt, split into ~250-word chunks."""
    chunks = []
    if not os.path.exists(path):
        print(f"  ⚠️  Not found: {path}")
        return chunks

    with open(path, "r", encoding="utf-8") as f:
        content = f.read()

    words = content.split()
    filename = os.path.basename(path)

    for i in range(0, len(words), WEBSITE_CHUNK_WORDS):
        chunk_words = words[i: i + WEBSITE_CHUNK_WORDS]
        if chunk_words:
            chunks.append({
                "text": " ".join(chunk_words),
                "source": filename,
                "source_type": "website",
                "language": "en",
                "chunk_index": i // WEBSITE_CHUNK_WORDS,
            })

    print(f"  ✅ {filename}: {len(chunks)} chunks extracted")
    return chunks


# ── Language-specific manual documents ───────────────────────────────────────
# Add .txt files directly to knowledge-base/ with a language suffix in filename:
#   product_data_tamil.txt  → language "ta"
#   product_data_hindi.txt  → language "hi"
#   product_data_malayalam.txt → language "ml"
#   product_data_english.txt   → language "en"
# These are chunked the same way as website data.

LANG_SUFFIX_MAP = {
    "tamil": "ta", "hindi": "hi", "malayalam": "ml",
    "english": "en", "kannada": "kn", "telugu": "te",
}


def detect_lang_from_filename(filename: str) -> str:
    """Guess language from filename suffix, e.g. product_data_tamil.txt → ta."""
    name = filename.lower().replace(".txt", "")
    for suffix, code in LANG_SUFFIX_MAP.items():
        if name.endswith(f"_{suffix}") or name.endswith(f"-{suffix}"):
            return code
    return "en"


def extract_manual_doc(path: str) -> list[dict]:
    """Chunk a manually added language-specific .txt from knowledge-base/."""
    chunks = []
    filename = os.path.basename(path)
    language = detect_lang_from_filename(filename)

    with open(path, "r", encoding="utf-8") as f:
        content = f.read()

    words = content.split()
    for i in range(0, len(words), WEBSITE_CHUNK_WORDS):
        chunk_words = words[i: i + WEBSITE_CHUNK_WORDS]
        if chunk_words:
            chunks.append({
                "text": " ".join(chunk_words),
                "source": filename,
                "source_type": "manual",
                "language": language,
                "chunk_index": i // WEBSITE_CHUNK_WORDS,
            })

    print(f"  ✅ {filename} [{language}]: {len(chunks)} chunks extracted")
    return chunks


# ── Main ──────────────────────────────────────────────────────────────────────

def process_all_knowledge() -> list[dict]:
    """Process all PPTX files, website data, and manual language docs."""
    all_chunks = []

    print("\n📄 Processing PPTX files...")
    for language, filename in PPTX_FILES.items():
        path = os.path.join(KNOWLEDGE_BASE_DIR, filename)
        chunks = extract_text_from_pptx(path, language)
        all_chunks.extend(chunks)

    print("\n🌐 Processing scraped website data...")
    if os.path.isdir(SCRAPED_DIR):
        for fname in sorted(os.listdir(SCRAPED_DIR)):
            if fname.endswith(".txt"):
                path = os.path.join(SCRAPED_DIR, fname)
                chunks = extract_text_from_scraped_website(path)
                all_chunks.extend(chunks)
    else:
        print(f"  ⚠️  Scraped data directory not found: {SCRAPED_DIR}")

    print("\n📝 Processing manual language documents...")
    found_manual = False
    for fname in sorted(os.listdir(KNOWLEDGE_BASE_DIR)):
        if fname.endswith(".txt"):
            found_manual = True
            path = os.path.join(KNOWLEDGE_BASE_DIR, fname)
            chunks = extract_manual_doc(path)
            all_chunks.extend(chunks)
    if not found_manual:
        print("  (none found — add .txt files to knowledge-base/ for language-specific data)")

    return all_chunks


def main():
    os.makedirs(CHUNKS_DIR, exist_ok=True)

    chunks = process_all_knowledge()

    if not chunks:
        print("\n❌ No chunks generated. Check that PPTX files exist in knowledge-base/")
        return

    with open(OUTPUT_FILE, "w", encoding="utf-8") as f:
        json.dump(chunks, f, ensure_ascii=False, indent=2)

    # Summary
    by_language: dict[str, int] = {}
    by_source_type: dict[str, int] = {}
    for c in chunks:
        lang = c.get("language", "unknown")
        stype = c.get("source_type", "unknown")
        by_language[lang] = by_language.get(lang, 0) + 1
        by_source_type[stype] = by_source_type.get(stype, 0) + 1

    print(f"\n✅ Total chunks: {len(chunks)}")
    print("   By language:", by_language)
    print("   By source type:", by_source_type)
    print(f"   Saved to: {OUTPUT_FILE}")


if __name__ == "__main__":
    main()
